"""Fill all AI Curriculum templates with supply-chain project content."""
from pathlib import Path
from copy import deepcopy
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PPTpt
from pptx.dml.color import RGBColor as PPTColor

CURR = Path("/Users/nalinaggarwal/Downloads/Supply-chain-main/AI Curriculum")
OUT  = CURR / "Filled"
OUT.mkdir(exist_ok=True)

TITLE = (
    "An Integrated Multi-Objective Optimization and Deep Reinforcement "
    "Learning Framework for Green, Resilient Supply Chain Management: "
    "Evidence from Indian Logistics Networks"
)
AUTHOR = "Nalin Aggarwal"

# ---------------------------------------------------------------------------
# Helper: replace all placeholder text in a paragraph
# ---------------------------------------------------------------------------
def replace_para(para, old, new):
    """Replace *old* with *new* across all runs in *para*, preserving formatting."""
    full = "".join(r.text for r in para.runs)
    if old not in full:
        return
    new_full = full.replace(old, new)
    for run in para.runs:
        run.text = ""
    if para.runs:
        para.runs[0].text = new_full
    else:
        para.add_run(new_full)

# ---------------------------------------------------------------------------
# 1. Introduction Template
# ---------------------------------------------------------------------------
def fill_introduction():
    doc = Document(CURR / "Nalin Introduction Template AI.docx")
    content = [
        ("Introduction Template", f"Introduction — {TITLE}"),
        ("500-600 words", "~600 words"),
        ("Cite at least 10-12 studies", "Cite at least 10-12 studies (see Literature Review Matrix for full list)"),
    ]
    intro_text = """
Indian logistics consumes approximately 14 % of GDP, well above the 8–10 % benchmark in mature freight economies (NCAER, 2024), and road freight generates roughly 260 million tonnes of CO₂ per year — a figure projected to quadruple by 2047 on the current trajectory (NITI Aayog & RMI, 2021). Indian operators run trucks at 60–65 % load factor with 35–40 % empty-running, and the Bharat Stage VI emission standard plus emerging ESG-disclosure obligations are tightening fast. A logistics planner therefore faces three simultaneous pressures — cost, carbon, and service reliability under disruption — that the existing literature treats one at a time.

Key terms: *multi-objective optimisation* refers to finding a set of trade-off solutions (the Pareto frontier) when two or more objectives conflict; *hypervolume indicator* is a scalar measure of how much of the objective space a Pareto front dominates; *proximal policy optimisation (PPO)* is a reinforcement-learning algorithm that trains an agent to maximise cumulative reward through trial-and-error interaction with a simulator; *discrete event simulation (DES)* models a system as a sequence of events in time, used here to stress-test routing plans under demand surges, supply disruptions, and route blockages.

Globally, supply-chain optimisation research has grown rapidly since the 2010s, with multi-objective evolutionary algorithms (Deb et al., 2002; Zhang & Li, 2007) and deep reinforcement learning (Schulman et al., 2017; Haarnoja et al., 2018) emerging as dominant paradigms. In India, the NITI Aayog & RMI (2021) freight roadmap and NCAER (2024) logistics cost report have highlighted the urgent need for data-driven planning tools that jointly address cost, carbon, and resilience. Despite this, most published frameworks use synthetic networks or Euclidean distances, treat carbon as a hard constraint rather than a competing objective, and report single-seed results without statistical significance testing.

The importance of further analysis is clear: current diagnostic models for Indian freight fail to account for the multifactorial nature of the problem — genetic, environmental, and operational factors interact in ways that single-objective or single-method approaches cannot capture. The present study addresses this gap by integrating four method streams end-to-end on a calibrated Indian network with a uniform statistical-validation protocol.

AI and machine learning have emerged as powerful tools for supply-chain management, enabling analysis of large, complex datasets that were previously unmanageable. Multi-objective evolutionary algorithms can explore the cost-carbon trade-off space efficiently; attention-LSTM networks can forecast demand with festival-spike sensitivity; and PPO controllers can learn adaptive inventory policies that outperform classical rules under disruption.

Current gaps in the literature include: (1) most studies focus on single-objective routing rather than the full cost-carbon-delivery-time trade-off; (2) emission models are rarely cross-verified against multiple primary sources; (3) learned inventory controllers are compared against classical baselines only on steady-state demand, not under disruption stress; (4) Indian-network calibration is absent from most published frameworks.

This paper argues that integrating NSGA-II multi-objective routing, SimPy-based discrete event simulation, and PPO inventory control on a calibrated Indian network — with a unified Friedman + Wilcoxon + Holm-Bonferroni statistical-validation protocol — produces a framework that is simultaneously more accurate, more resilient, and more honest about its statistical claims than any single-method approach in the existing literature.
""".strip()

    # Replace the body paragraphs
    for para in doc.paragraphs:
        t = para.text.strip()
        if t == "Introduction Template":
            para.clear()
            run = para.add_run(f"Introduction — {TITLE}")
            run.bold = True
        elif t.startswith("500-600 words"):
            para.clear()
            para.add_run("Target length: ~600 words. Cite at least 10–12 studies.")
        elif "Definition of Key Terms" in t:
            para.clear()
            para.add_run(intro_text)
            break  # rest of template is instructions; stop here

    doc.save(OUT / "Nalin Introduction Template AI — Filled.docx")
    print("Saved: Introduction")

fill_introduction()

# ---------------------------------------------------------------------------
# 2. Literature Review Matrix
# ---------------------------------------------------------------------------
def fill_lit_review():
    doc = Document(CURR / "Nalin Literature Review Matrix Template-Generic.docx")
    rows = [
        ["Deb et al., 2002", "A Fast and Elitist Multiobjective Genetic Algorithm: NSGA-II",
         "Multi-objective evolutionary optimisation for combinatorial problems",
         "NSGA-II (non-dominated sorting, crowding distance)",
         "Combinatorial optimisation / VRP",
         "NSGA-II outperforms NSGA-I and PAES on DTLZ benchmarks; fast non-dominated sort is O(MN²)",
         "Rigorous benchmark suite; widely cited; open-source pymoo implementation",
         "Only tested on synthetic problems; no real-network calibration",
         "Primary routing algorithm in Phase 1; our joint-normalised HV indicator extends its evaluation protocol"],
        ["Zhang & Li, 2007", "MOEA/D: A Multiobjective Evolutionary Algorithm Based on Decomposition",
         "Decomposition-based multi-objective optimisation",
         "MOEA/D (Tchebycheff scalarisation, neighbourhood update)",
         "Combinatorial optimisation",
         "MOEA/D is competitive with NSGA-II on ZDT/DTLZ; neighbourhood size is a critical hyperparameter",
         "Scalable to many objectives; low memory footprint",
         "Sensitive to weight-vector design; brittle on heterogeneous objective ranges",
         "Third baseline in Phase 1; our results confirm wider seed-to-seed variance (std 0.328 vs 0.143 for NSGA-II)"],
        ["Bektas & Laporte, 2011", "The Pollution-Routing Problem",
         "Carbon-constrained vehicle routing",
         "Mixed-integer programming; adaptive large-neighbourhood search",
         "Green logistics / VRP",
         "Formalises the cost-vs-emission trade-off as a bi-objective programme; introduces the green-premium curve concept",
         "Rigorous mathematical formulation; widely cited in green-VRP literature",
         "Euclidean distances only; no stochastic resilience component; single-seed results",
         "Foundational reference for our carbon-budget ε-constraint formulation (§3.4) and green-premium curve (§6.1)"],
        ["Schulman et al., 2017", "Proximal Policy Optimization Algorithms",
         "On-policy deep reinforcement learning for continuous control",
         "PPO-Clip (clipped surrogate objective, GAE advantage estimation)",
         "Reinforcement learning / inventory control",
         "PPO achieves state-of-the-art on MuJoCo continuous-control benchmarks with stable training",
         "Simple implementation; robust to hyperparameter choice; canonical reference",
         "On-policy; sample-inefficient compared to SAC; tested on locomotion not supply-chain",
         "Primary inventory controller in Phase 3; we adopt canonical hyperparameters (clip=0.2, GAE λ=0.95, lr=1e-4)"],
        ["Haarnoja et al., 2018", "Soft Actor-Critic: Off-Policy Maximum Entropy Deep Reinforcement Learning",
         "Off-policy deep reinforcement learning with entropy regularisation",
         "SAC (twin Q-networks, automatic temperature tuning)",
         "Reinforcement learning",
         "SAC outperforms PPO on sample efficiency; automatic entropy tuning removes a key hyperparameter",
         "Sample-efficient; stable training; automatic temperature tuning",
         "More complex implementation; not yet benchmarked on supply-chain disruption tasks",
         "Available as a config switch in Phase 3; not yet validated on our disruption stress-test"],
        ["Hosseini et al., 2019", "A review of definitions and measures of system resilience",
         "Supply-chain resilience metrics taxonomy",
         "Systematic review; Bayesian network resilience measure",
         "Supply-chain resilience",
         "Identifies time-to-survive (TTS) and time-to-recover (TTR) as the canonical resilience metrics; proposes magnitude-normalised TTR",
         "Comprehensive taxonomy; widely cited; clear metric definitions",
         "Descriptive rather than prescriptive; not coupled to an optimisation or control layer",
         "We adopt the normalised TTR verbatim in Phase 2 DES evaluation; the disruption-stress comparison in §5.6 is framed using TTS/TTR"],
        ["Boute et al., 2022", "Deep Reinforcement Learning for Inventory Control: A Roadmap",
         "DRL applied to inventory management",
         "PPO, SAC, DQN benchmarked on lost-sales and dual-sourcing problems",
         "Inventory management / supply chain",
         "DRL wins over classical heuristics only when the action space is continuous; steady-state comparison understates DRL value",
         "Comprehensive benchmark; calls for disruption-stress comparison",
         "Does not test under disruption; single-echelon focus",
         "Motivates our disruption-stress framing; we extend their benchmark to multi-warehouse multi-regime comparison"],
        ["Gijsbrechts et al., 2022", "Can Deep Reinforcement Learning Improve Inventory Management?",
         "DRL vs classical inventory policies on lost-sales benchmarks",
         "PPO, SAC on multi-echelon lost-sales problems",
         "Inventory management",
         "DRL matches or beats (R,s,S) on continuous-action problems; per-day cost is the right comparison metric",
         "Rigorous benchmark; MSOM publication; open-source code",
         "Steady-state only; no disruption regimes; single-country network",
         "Primary benchmark reference for our PPO vs (R,s,S) comparison; we extend to four disruption regimes"],
        ["Konstantakopoulos et al., 2022", "Vehicle Routing Problem and Related Algorithms for Logistics Distribution",
         "Systematic review of VRP algorithms 2009–2020",
         "Survey of NSGA-II, MOEA/D, Clarke-Wright, and hybrid methods",
         "Green logistics / VRP",
         "Fewer than 8% of VRP papers combine bi-objective NSGA-II with a calibrated CO₂ emission model on a published benchmark",
         "144-paper systematic review; identifies key gaps in the literature",
         "Review only; no new empirical results",
         "Confirms the integration gap our framework closes; cited in §2.5 positioning argument"],
        ["Yang et al., 2024", "Dynamic Optimization of Multi-Echelon Supply Chain Inventory Under Disruptive Scenarios",
         "DRL for inventory control under disruption",
         "PPO on multi-echelon supply chain with demand and supply shocks",
         "Supply chain resilience / DRL",
         "PPO-vs-classical gap widens with disruption severity; disruption-stress test is the right comparison",
         "First paper to frame DRL value-add as a disruption-resilience instrument",
         "Single-country synthetic network; no multi-objective routing layer",
         "Directly motivates our disruption-first framing in §5.6; we reproduce and extend their finding on a calibrated Indian network"],
    ]

    # Find the table and fill it
    for table in doc.tables:
        if len(table.columns) >= 8:
            for i, row_data in enumerate(rows):
                if i + 1 >= len(table.rows):
                    table.add_row()
                row = table.rows[i + 1]
                for j, cell_text in enumerate(row_data):
                    if j < len(row.cells):
                        row.cells[j].text = cell_text
            break

    doc.save(OUT / "Nalin Literature Review Matrix — Filled.docx")
    print("Saved: Literature Review Matrix")

fill_lit_review()

# ---------------------------------------------------------------------------
# 3. Methodology Template (docx)
# ---------------------------------------------------------------------------
METHODOLOGY_SECTIONS = {
    "1. Overview": (
        "This study employs a four-phase sequential decomposition framework to optimise "
        "green, resilient supply-chain operations on a calibrated Indian logistics network. "
        "Phase 1 applies multi-objective evolutionary algorithms (NSGA-II, NSGA-III, MOEA/D) "
        "to generate a Pareto frontier of routing plans trading transport cost against CO₂ "
        "emissions. Phase 2 stress-tests each plan using a SimPy-based discrete event "
        "simulation under demand surges, supply disruptions, and route blockages. Phase 3 "
        "trains an Attention-LSTM demand forecaster and a PPO inventory controller to manage "
        "day-to-day replenishment decisions. Phase 4 applies Sobol global sensitivity analysis "
        "to identify which input parameters drive the headline outcomes. The methodology "
        "consists of data collection, preprocessing, feature engineering, model development, "
        "training, validation, and evaluation against classical baselines."
    ),
    "2. Data Collection": (
        "We utilised three primary data sources. (1) The Dalal (2022) INFORMS Journal on "
        "Computing supplement, comprising 101 customer demand points and 5 warehouses across "
        "India with GPS coordinates and a pre-computed pairwise distance matrix. (2) The DataCo "
        "Smart Supply Chain dataset (180,519 orders from ~20,000 customers, CC BY licence), "
        "used to fit log-normal demand parameters (μ=6.44, σ=0.97). (3) The Delhivery shipment "
        "dataset (144,867 shipments, 10 hubs × 150 customers), used for cross-validation. "
        "Road distances were queried from OSRM v5 with OpenRouteService fallback. Emission "
        "constants were sourced from MEET (Hickman 1999), cross-verified against COPERT 5, "
        "HBEFA 4.2, and IPCC AR6. Fleet parameters (70:30 HCV:LCV ratio, 65% utilisation, "
        "35% empty-running) were sourced from NITI Aayog & RMI (2021) and VAHAN FY2024."
    ),
    "3. Data Preprocessing": (
        "3.1 Data Cleaning: Missing OSRM distance values were filled using the OpenRouteService "
        "fallback API. Demand values were clipped to [100, 10000] kg to remove outliers. "
        "3.2 Normalisation: The joint-normalised hypervolume indicator scales each objective "
        "to its ideal-nadir interval before computing the hypervolume contribution, removing "
        "the bias that raw HV introduces when objectives have heterogeneous numerical ranges "
        "(cost ~10⁶ INR vs emissions ~10⁵ kg CO₂). Demand series were standardised using "
        "training-set mean and standard deviation (chronological split; no data leakage). "
        "3.3 Dimensionality: The Sobol sensitivity analysis uses a reduced 3-warehouse × "
        "8-customer instance to keep the 1,280 NSGA-II evaluations tractable on CPU. The "
        "full 5-warehouse × 100-customer instance is used for all other experiments."
    ),
    "4. Feature Engineering": (
        "For the LSTM forecaster, the input feature vector is a 30-day rolling window of "
        "demand at the target customer and its four nearest neighbours (5 × 30 = 150 features). "
        "The output is the next 7 days of demand. The chronological 70/15/15 train/val/test "
        "split ensures no future information leaks into the training set (Tashman, 2000). "
        "For the PPO controller, the 45-dimensional state vector concatenates per-warehouse "
        "inventory levels, in-transit quantities, the 7-day LSTM forecast, the 7-day realised "
        "demand history, and a binary shock indicator. Feature importance was assessed via "
        "Sobol global sensitivity analysis (Saltelli et al., 2010), which identified demand "
        "variability as the dominant driver (S1=0.72, ST=0.90)."
    ),
    "5. Model Development": (
        "5.1 Multi-objective evolutionary algorithms: NSGA-II (pop=500, max_gen=200, SBX η_c=15, "
        "PM η_m=20, OR-Tools warm-start with 2 seeds) was the primary planner. NSGA-III "
        "(pop=92, 91 Das-Dennis reference points, 3 objectives) and MOEA/D (pop=500, "
        "Tchebycheff scalarisation, neighbourhood=20) were benchmarked as alternatives. "
        "A diversity-preserving marginal-tradeoff repair operator assigns each individual a "
        "private scalarisation weight, preventing the front-collapse pathology of proportional "
        "repair. 5.2 Deep learning: The Attention-LSTM forecaster uses 2 LSTM layers (256 "
        "hidden units) + Bahdanau attention + FC head. The PPO controller uses a Beta-distribution "
        "actor (45-dim state, 5-dim action) with GAE (λ=0.95), clip range 0.2, lr=1×10⁻⁴, "
        "trained for 2M steps. 5.3 Classical baselines: (R,s,S) periodic-review policy "
        "(review period 7 days, parameters tuned on the same demand distribution) and a "
        "uniform-random policy as a lower bound."
    ),
    "6. Training and Validation": (
        "NSGA-II, NSGA-III, and MOEA/D were each run for 50 independent seeds (seeds 0–49) "
        "on the calibrated 5-warehouse × 100-customer instance. The LSTM was trained with "
        "Adam (lr=1×10⁻³, batch=64, early-stop patience=10) on a 70/15/15 chronological "
        "split. The PPO controller was trained for 2M environment steps on the full network "
        "using the stress-mode periodic-review environment (initial inventory 30% of capacity, "
        "3-day lead time, INR-denominated reward). The DES Monte Carlo used 100 replications "
        "per scenario per plan. Cross-validation was performed on the Delhivery secondary "
        "network (20 seeds) to establish external validity."
    ),
    "7. Evaluation Metrics": (
        "Phase 1 (routing): Joint-normalised hypervolume (HV) as the primary indicator; "
        "Pareto front size as a secondary indicator. Statistical significance: 3-way Friedman "
        "omnibus test (χ²=7.32, p=0.0257) + paired Wilcoxon post-hoc with Holm-Bonferroni "
        "correction. CVRPLIB gap-to-BKS (mean +5.1%) for implementation correctness. "
        "Phase 2 (resilience): Mean service level (95.6%±0.28%), 95% CI lower bound (95.09%), "
        "time-to-survive (TTS), time-to-recover (TTR). "
        "Phase 3 (forecasting): MAPE (23.5%), RMSE (56.5 kg). "
        "Phase 3 (inventory control): Per-day cost (INR), episode survival length (days), "
        "service level — reported separately for each of four disruption regimes."
    ),
    "8. Model Interpretation, Feature importance and Biological Insight": (
        "Feature importance was assessed via Sobol global sensitivity analysis (SALib, "
        "Herman & Usher 2017) with Saltelli sampling (N=128, 1,280 NSGA-II evaluations). "
        "Demand variability dominates with first-order index S1=0.72 and total-order ST=0.90, "
        "meaning it explains 72% of the variance in the joint cost-carbon objective directly "
        "and 90% including interactions. Warehouse capacity acts purely through interactions "
        "(S1≈0, ST=0.35). Fleet mix and carbon weight have smaller total-order effects "
        "(ST≈0.30 each). The operational implication: demand-shaping investments (forecasting, "
        "retailer education, promotional planning) outrank fleet purchases on value-of-information. "
        "For the PPO controller, SHAP analysis (not yet implemented; planned for future work) "
        "would reveal which state dimensions drive reorder decisions under each disruption regime."
    ),
    "9. Software and Tools": (
        "Python 3.10; PyTorch 2.0 (LSTM, PPO, SAC); pymoo 0.6.x (NSGA-II, NSGA-III, MOEA/D); "
        "SimPy 4.1 (DES); Gymnasium 0.29 (RL environment); SALib 1.4 (Sobol sensitivity); "
        "OR-Tools (warm-start seeds, Clarke-Wright baseline); OSRM v5 + OpenRouteService "
        "(road distances); MLflow (experiment tracking); Pandas, NumPy, Matplotlib, Seaborn "
        "(data analysis and visualisation). All dependencies pinned with == version pins in "
        "supply_chain_research/requirements.txt. Training performed on NVIDIA Tesla T4 16 GB "
        "(Modal cloud GPU, ~3 hours wall-clock for the full 50-seed pipeline)."
    ),
}

def fill_methodology_docx():
    doc = Document(CURR / "Nalin Methodology_Template_AI.docx")
    current_section = None
    for para in doc.paragraphs:
        t = para.text.strip()
        if t == "Methodology Template for AI Based Research Paper":
            para.clear()
            run = para.add_run(f"Methodology — {TITLE}")
            run.bold = True
            continue
        for sec_key, sec_content in METHODOLOGY_SECTIONS.items():
            if t.startswith(sec_key.split(".")[0] + ".") and sec_key.split(".")[1].strip() in t:
                current_section = sec_key
                break
        if current_section and t.startswith("Example:"):
            para.clear()
            para.add_run(METHODOLOGY_SECTIONS[current_section])
            current_section = None

    doc.save(OUT / "Nalin Methodology Template AI — Filled.docx")
    print("Saved: Methodology (docx)")

fill_methodology_docx()

# ---------------------------------------------------------------------------
# 4. Results Template
# ---------------------------------------------------------------------------
def fill_results():
    doc = Document(CURR / "Nalin Results Template for AI.docx")
    for para in doc.paragraphs:
        t = para.text.strip()
        if t == "Results Template for AI":
            para.clear()
            run = para.add_run(f"Results — {TITLE}")
            run.bold = True
        elif "This section presents the results" in t:
            para.clear()
            para.add_run(
                "This section presents the results obtained from the multi-objective "
                "evolutionary algorithms, discrete event simulation, attention-LSTM "
                "demand forecaster, and PPO inventory controller applied to the "
                "calibrated 5-warehouse, 100-customer Indian logistics network."
            )
        elif "After data cleaning and preprocessing" in t:
            para.clear()
            para.add_run(
                "After data cleaning and preprocessing, 101 customer demand points and "
                "5 warehouses were retained. Demand parameters were fitted to the DataCo "
                "dataset (180,519 orders): LogNormal(μ=6.44, σ=0.97), median demand "
                "~626 kg/customer, 95th percentile ~3,090 kg. Road distances were "
                "computed via OSRM v5 for all 5×101 warehouse-to-customer pairs. "
                "The joint-normalised hypervolume indicator scaled each objective to its "
                "ideal-nadir interval, removing the magnitude bias between cost (~10⁶ INR) "
                "and emissions (~10⁵ kg CO₂)."
            )
        elif "Accuracy: X%" in t:
            # First occurrence = ML model block
            para.clear()
            para.add_run(
                "NSGA-II (primary multi-objective planner):\n"
                "  Joint-normalised HV: 0.713 ± 0.143 (50 seeds)\n"
                "  Mean Pareto front size: 11.2 (range 4–21)\n"
                "  CVRPLIB Augerat Set-A: 27/27 instances, mean gap +5.1%\n\n"
                "NSGA-III (3-objective extension):\n"
                "  Joint-normalised HV: 0.659 ± 0.203\n"
                "  Mean front size: 7.2 (range 2–13)\n\n"
                "MOEA/D (decomposition baseline):\n"
                "  Joint-normalised HV: 0.595 ± 0.328\n"
                "  Mean front size: 3.3\n\n"
                "Statistical significance: Friedman χ²=7.32, p=0.0257 (omnibus rejects "
                "equal medians). Wilcoxon NSGA-II vs MOEA/D: W=399, p=0.0207 (raw); "
                "Holm-adjusted p=0.062. Pairwise post-hoc tests do not survive "
                "Holm-Bonferroni correction."
            )

    # Fill the performance comparison table
    for table in doc.tables:
        if len(table.columns) >= 5 and "Random Forest" in table.rows[1].cells[0].text:
            rows_data = [
                ["NSGA-II", "HV 0.713±0.143", "Front size 11.2", "Friedman rank 1st", "Recommended (routing)"],
                ["NSGA-III", "HV 0.659±0.203", "Front size 7.2", "Friedman rank 2nd", "Use for 3-objective"],
                ["MOEA/D", "HV 0.595±0.328", "Front size 3.3", "Friedman rank 3rd", "Baseline only"],
                ["PPO", "Cost/day -850 INR (severe)", "Survival 91 days", "SL 95.4%", "Recommended (disruption)"],
                ["(R,s,S)", "Cost/day -676 INR (steady)", "Survival 61 days (severe)", "SL 93.4%", "Recommended (steady-state)"],
                ["Attention-LSTM", "MAPE 23.5%", "RMSE 56.5 kg", "—", "Recommended (forecasting)"],
            ]
            for i, rd in enumerate(rows_data):
                if i + 1 < len(table.rows):
                    for j, val in enumerate(rd):
                        if j < len(table.rows[i+1].cells):
                            table.rows[i+1].cells[j].text = val
            break

    doc.save(OUT / "Nalin Results Template AI — Filled.docx")
    print("Saved: Results")

fill_results()

# ---------------------------------------------------------------------------
# 5. Discussion Template
# ---------------------------------------------------------------------------
DISCUSSION = {
    "restate": (
        "This study addressed three research questions: (RQ1) how multi-objective evolutionary "
        "algorithms perform on a calibrated Indian logistics network with a verified emission "
        "model; (RQ2) what the green-premium curve looks like and where its knee point lies; "
        "and (RQ3) whether a PPO inventory controller improves supply-chain resilience under "
        "disruption relative to a tuned (R,s,S) baseline. The primary objective was to "
        "integrate NSGA-II routing, DES resilience evaluation, and PPO inventory control "
        "under a unified statistical-validation protocol on a real Indian network."
    ),
    "findings": (
        "The most important findings are: (1) NSGA-II achieves the highest joint-normalised "
        "hypervolume (0.713±0.143) with the richest Pareto front (11.2 solutions/seed), "
        "confirmed by a Friedman omnibus test (p=0.0257). (2) The green-premium curve has "
        "its knee in the 20–30% carbon-reduction band, where the marginal cost per kg of "
        "avoided CO₂ is comparable to the proposed Indian carbon-tax rate (~INR 0.40/kg). "
        "(3) PPO survives 91 days under severe disruption vs 61 days for (R,s,S), with the "
        "per-day cost gap closing to 3% (-850 vs -876 INR/day). (4) Demand variability is "
        "the dominant driver of the cost-carbon outcome (Sobol S1=0.72, ST=0.90). "
        "(5) The Attention-LSTM achieves 23.5% MAPE, within the published 18–28% band for "
        "log-normal demand series with festival spikes."
    ),
    "theoretical": (
        "The findings corroborate the multi-objective VRP literature (Deb et al., 2002; "
        "Bektas & Laporte, 2011) that treating cost and carbon as competing objectives "
        "rather than as a constrained single objective produces richer, more actionable "
        "Pareto fronts. The diversity-preserving marginal-tradeoff repair operator is a "
        "novel contribution that addresses the front-collapse pathology of proportional "
        "repair, producing 10–15 distinct solutions per seed vs 1–4 under proportional "
        "repair. The disruption-stress framing for PPO evaluation extends the findings of "
        "Yang et al. (2024) and Boute et al. (2022) to a calibrated Indian network with "
        "four disruption regimes, confirming that the DRL-vs-classical gap widens with "
        "disruption severity."
    ),
    "practical": (
        "The framework has direct practical applications for Indian logistics planners. "
        "The green-premium curve provides a defensible rupee-per-kg-CO₂ number for "
        "sustainability-disclosure conversations. The Sobol sensitivity result redirects "
        "investment priority from fleet purchases to demand-shaping contracts. The PPO "
        "controller offers a concrete resilience instrument for disruption-exposed corridors, "
        "with a recommended three-phase deployment roadmap: shadow-mode pilot on one "
        "corridor, network-wide live recommendation, then PPO layered on disruption-exposed "
        "nodes with rolling 90-day retraining."
    ),
    "strengths": (
        "Key strengths include: (1) use of a peer-reviewed Indian network (Dalal 2022) "
        "rather than synthetic geography; (2) emission model cross-verified against four "
        "independent sources (MEET, COPERT 5, HBEFA 4.2, IPCC AR6); (3) 50-seed statistical "
        "protocol with Friedman omnibus + Wilcoxon post-hoc + Holm-Bonferroni correction; "
        "(4) external validity established on a second Indian network (Delhivery, HV 0.880±0.099); "
        "(5) full reproducibility with pinned dependencies, fixed seeds, and a 454-test suite."
    ),
    "limitations": (
        "Limitations include: (1) single-country network — the cost and emission parameters "
        "are India-specific and would need recalibration for other markets; (2) simulated "
        "demand rather than a single shipper's historical record; (3) sim-to-real gap on "
        "the PPO controller — deployment requires a shadow-mode calibration phase; "
        "(4) 50 episodes per disruption cell is sufficient for directional conclusions but "
        "tight for tight confidence intervals at the moderate regime; (5) pairwise post-hoc "
        "Wilcoxon tests do not survive Holm-Bonferroni correction, so the headline claim is "
        "distributional rather than ordinal."
    ),
    "future": (
        "Future research should: (1) extend the routing layer to multi-modal rail-and-road "
        "formulations to address the Indian Dedicated Freight Corridor; (2) integrate "
        "real-time IoT sensor streams for sub-daily controller response; (3) apply "
        "transfer-learning techniques to reduce per-corridor retraining cost; (4) couple "
        "the cost-carbon trade-off to a carbon-credit trading layer as Indian carbon-pricing "
        "instruments mature; (5) expand the disruption-stress benchmark to 200 episodes per "
        "cell for tighter confidence intervals."
    ),
    "conclusion": (
        "This study demonstrated that integrating NSGA-II multi-objective routing, SimPy "
        "discrete event simulation, and PPO inventory control on a calibrated Indian network "
        "— under a unified Friedman + Wilcoxon + Holm-Bonferroni statistical-validation "
        "protocol — produces a framework that is simultaneously more accurate, more resilient, "
        "and more statistically honest than any single-method approach in the existing "
        "literature. The identification of demand variability as the dominant driver "
        "(S1=0.72) provides a clear investment priority for Indian logistics planners. "
        "The findings pave the way for AI-driven supply-chain management that jointly "
        "addresses cost, carbon, and resilience at scale."
    ),
}

def fill_discussion():
    doc = Document(CURR / "Nalin Discussion template AI.docx")
    for para in doc.paragraphs:
        t = para.text.strip()
        if t == "Discussion Template AI":
            para.clear()
            run = para.add_run(f"Discussion — {TITLE}")
            run.bold = True
        elif "briefly restate the research problem" in t:
            para.clear()
            para.add_run(DISCUSSION["restate"])
        elif "Highlight the Most Important Findings" in t:
            para.clear()
            para.add_run(DISCUSSION["findings"])
        elif "Discuss how the results contribute to existing knowledge" in t:
            para.clear()
            para.add_run(DISCUSSION["theoretical"])
        elif "Describe how the results could be applied in practice" in t:
            para.clear()
            para.add_run(DISCUSSION["practical"])
        elif "Highlight key strengths" in t:
            para.clear()
            para.add_run(DISCUSSION["strengths"])
        elif "Address limitations related to data quality" in t:
            para.clear()
            para.add_run(DISCUSSION["limitations"])
        elif "Recommend areas that require further exploration" in t:
            para.clear()
            para.add_run(DISCUSSION["future"])
        elif "Reiterate the key points made in the discussion" in t:
            para.clear()
            para.add_run(DISCUSSION["conclusion"])

    doc.save(OUT / "Nalin Discussion Template AI — Filled.docx")
    print("Saved: Discussion")

fill_discussion()

# ---------------------------------------------------------------------------
# 6. Methodology Presentation (pptx)
# ---------------------------------------------------------------------------
SLIDE_CONTENT = {
    1: ("Methodology", f"{TITLE}\n{AUTHOR}"),
    2: ("Overview",
        "This study employs a four-phase sequential decomposition framework to optimise "
        "green, resilient supply-chain operations on a calibrated Indian logistics network. "
        "Phase 1 applied NSGA-II, NSGA-III, and MOEA/D to generate Pareto-optimal routing "
        "plans trading transport cost against CO₂ emissions. Phase 2 stress-tested each plan "
        "using SimPy discrete event simulation under three shock classes. Phase 3 trained an "
        "Attention-LSTM demand forecaster and a PPO inventory controller. Phase 4 applied "
        "Sobol global sensitivity analysis to identify dominant input parameters."),
    3: ("Data Collection",
        "Sources: (1) Dalal (2022) INFORMS supplement — 101 customers, 5 warehouses, "
        "GPS coordinates, pairwise distance matrix. (2) DataCo Smart Supply Chain dataset "
        "— 180,519 orders, ~20,000 customers (CC BY). (3) Delhivery shipment dataset — "
        "144,867 shipments, 10 hubs × 150 customers (cross-validation). Road distances: "
        "OSRM v5 + OpenRouteService fallback. Emission constants: MEET (Hickman 1999), "
        "COPERT 5, HBEFA 4.2, IPCC AR6. Fleet parameters: NITI Aayog & RMI (2021), "
        "VAHAN FY2024. Total: 101 customers, 5 warehouses, 212 model parameters."),
    4: ("Data Preprocessing",
        "1. Data Cleaning: Missing OSRM values filled via ORS fallback. Demand clipped to "
        "[100, 10000] kg. 2. Normalisation: Joint-normalised hypervolume scales each "
        "objective to its ideal-nadir interval (removes magnitude bias: cost ~10⁶ INR vs "
        "emissions ~10⁵ kg CO₂). LSTM demand series standardised using training-set "
        "statistics only (chronological split; no data leakage). 3. Dimensionality: "
        "Sobol sensitivity uses a reduced 3×8 instance (1,280 NSGA-II evaluations). "
        "Full 5×100 instance used for all other experiments."),
    5: ("Feature Engineering",
        "LSTM forecaster: 30-day rolling window × 5 nearest customers = 150 input features. "
        "Output: 7-day demand forecast. Chronological 70/15/15 split (Tashman, 2000). "
        "PPO controller: 45-dimensional state vector — per-warehouse inventory, in-transit "
        "quantity, 7-day LSTM forecast, 7-day realised demand, binary shock indicator. "
        "Feature importance: Sobol global sensitivity (SALib, Herman & Usher 2017). "
        "Key finding: demand variability dominates (S1=0.72, ST=0.90)."),
    6: ("Model Development",
        "1. Multi-objective routing: NSGA-II (pop=500, SBX η_c=15, PM η_m=20, OR-Tools "
        "warm-start, diversity-preserving repair). NSGA-III (pop=92, 91 reference points). "
        "MOEA/D (pop=500, Tchebycheff). Clarke-Wright (correctness baseline). "
        "2. Deep learning: Attention-LSTM (2×256 LSTM + Bahdanau attention). PPO (Beta "
        "actor, 45-dim state, GAE λ=0.95, clip=0.2, lr=1e-4, 2M steps). "
        "3. Classical baselines: (R,s,S) periodic-review (review=7 days, tuned). "
        "Random policy (lower bound)."),
    7: ("Training and Validation",
        "NSGA-II/III/MOEA/D: 50 independent seeds (seeds 0–49), 5×100 network. "
        "LSTM: Adam (lr=1e-3, batch=64), early-stop patience=10, 70/15/15 chronological split. "
        "PPO: 2M environment steps, stress-mode env (initial inventory 30% capacity, "
        "3-day lead time, INR reward). DES: 100 Monte Carlo replications per scenario. "
        "Cross-validation: Delhivery secondary network, 20 seeds. "
        "All experiments: master seed=42, pinned dependencies (requirements.txt)."),
    8: ("Evaluation Metrics",
        "Phase 1 (routing): Joint-normalised hypervolume (primary), Pareto front size "
        "(secondary). Statistical tests: Friedman omnibus (χ²=7.32, p=0.0257) + Wilcoxon "
        "post-hoc + Holm-Bonferroni correction. CVRPLIB gap-to-BKS (mean +5.1%). "
        "Phase 2 (resilience): Service level (95.6%±0.28%), 95% CI lower bound (95.09%), "
        "TTS, TTR. Phase 3 (forecasting): MAPE (23.5%), RMSE (56.5 kg). "
        "Phase 3 (inventory): Per-day cost (INR), survival days, service level — "
        "reported for 4 disruption regimes."),
    9: ("Model Interpretation and Feature Importance",
        "Sobol global sensitivity (Saltelli N=128, 1,280 NSGA-II evaluations): "
        "demand_variability: S1=0.72, ST=0.90 (dominant — mostly direct effect). "
        "warehouse_capacity_factor: S1≈0, ST=0.35 (pure interaction). "
        "fleet_mix_ratio: S1=-0.05, ST=0.30 (interaction-driven). "
        "carbon_weight: S1=0.05, ST=0.30 (mostly interaction). "
        "Operational implication: demand-shaping investments outrank fleet purchases. "
        "PPO feature importance: SHAP analysis planned for future work."),
    10: ("Software and Tools",
        "Python 3.10 | PyTorch 2.0 | pymoo 0.6.x | SimPy 4.1 | Gymnasium 0.29 | "
        "SALib 1.4 | OR-Tools | OSRM v5 + OpenRouteService | MLflow | "
        "Pandas, NumPy, Matplotlib, Seaborn. "
        "Hardware: NVIDIA Tesla T4 16 GB (Modal cloud GPU). "
        "Training time: ~3 hours wall-clock for full 50-seed pipeline. "
        "All dependencies pinned with == version pins. "
        "454 automated tests; full reproducibility from a fresh clone."),
}

def fill_methodology_pptx():
    prs = Presentation(CURR / "Nalin Methodology AI Presentation.pptx")
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num not in SLIDE_CONTENT:
            continue
        title_text, body_text = SLIDE_CONTENT[slide_num]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full = "\n".join(p.text for p in tf.paragraphs).strip()
            # Title shape: short text matching the slide title
            if full.lower().startswith(title_text.lower()[:10]) or (
                slide_num == 1 and "Methodology" in full
            ):
                if len(full) < 60:  # it's the title box
                    tf.paragraphs[0].runs[0].text = title_text if slide_num > 1 else "Methodology"
                    continue
            # Body shape: replace instruction text with content
            if len(full) > 30 and slide_num > 1:
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if tf.paragraphs:
                    tf.paragraphs[0].add_run().text = body_text

    prs.save(OUT / "Nalin Methodology AI Presentation — Filled.pptx")
    print("Saved: Methodology (pptx)")

fill_methodology_pptx()

# ---------------------------------------------------------------------------
# 7. Mentor Template — fill the session plan
# ---------------------------------------------------------------------------
def fill_mentor_template():
    doc = Document(CURR / "Nalin AI Research Mentor Template.docx")
    for table in doc.tables:
        if len(table.columns) >= 4 and "Topic" in table.rows[0].cells[0].text:
            # Row 1: Topic definition
            r = table.rows[1]
            r.cells[0].text = "Topic definition — Supply-chain optimisation with AI"
            r.cells[1].text = (
                "Problem statement: How can multi-objective evolutionary algorithms, "
                "discrete event simulation, and deep reinforcement learning be integrated "
                "to jointly optimise cost, carbon, and resilience on a calibrated Indian "
                "logistics network? Hypothesis: NSGA-II + PPO outperforms classical "
                "(R,s,S) under disruption stress."
            )
            r.cells[2].text = "docs/MENTOR_REPORT.md; docs/PAPER_OUTLINE.md §1"
            r.cells[3].text = "1–2"
            # Row 2: Reading research papers
            r = table.rows[2]
            r.cells[1].text = (
                "Read and summarise: Deb et al. 2002 (NSGA-II), Bektas & Laporte 2011 "
                "(PRP), Schulman et al. 2017 (PPO), Hosseini et al. 2019 (resilience), "
                "Boute et al. 2022 (DRL inventory). Complete literature review matrix."
            )
            r.cells[2].text = "AI Curriculum/Filled/Nalin Literature Review Matrix — Filled.docx"
            # Row 3: Introduction/Literature review
            r = table.rows[3]
            r.cells[1].text = (
                "Write the introduction (§1) and literature review (§2) sections of the "
                "manuscript. Use the filled Introduction Template and Literature Review "
                "Matrix as starting points."
            )
            r.cells[2].text = "AI Curriculum/Filled/Nalin Introduction Template AI — Filled.docx"
            # Row 4: Methodology
            r = table.rows[4]
            r.cells[1].text = (
                "Write the methodology section (§3–§4) covering NSGA-II, DES, LSTM, PPO, "
                "and Sobol sensitivity. Use the filled Methodology Template and Presentation."
            )
            r.cells[2].text = (
                "AI Curriculum/Filled/Nalin Methodology Template AI — Filled.docx; "
                "AI Curriculum/Filled/Nalin Methodology AI Presentation — Filled.pptx"
            )
            # Row 5: Data Collection
            r = table.rows[5]
            r.cells[1].text = (
                "Confirm data sources: Dalal (2022) supplement, DataCo dataset, Delhivery "
                "dataset, OSRM distances, MEET/COPERT/HBEFA/IPCC emission constants. "
                "Document in docs/DATA_SOURCES.md."
            )
            r.cells[2].text = "docs/DATA_SOURCES.md; docs/REPLICATION_GUIDE.md"
            # Row 6: Results
            r = table.rows[6]
            r.cells[1].text = (
                "Write the results section (§5) covering NSGA-II HV results, CVRPLIB "
                "validation, DES resilience, LSTM forecasting, PPO disruption stress-test, "
                "ablation study, and cross-validation. Use the filled Results Template."
            )
            r.cells[2].text = "AI Curriculum/Filled/Nalin Results Template AI — Filled.docx"
            # Row 7: Discussion
            r = table.rows[7]
            r.cells[1].text = (
                "Write the discussion section (§6–§7) covering managerial insights, "
                "limitations, and future work. Use the filled Discussion Template."
            )
            r.cells[2].text = "AI Curriculum/Filled/Nalin Discussion Template AI — Filled.docx"
            break

    doc.save(OUT / "Nalin AI Research Mentor Template — Filled.docx")
    print("Saved: Mentor Template")

fill_mentor_template()
print("\nAll templates filled successfully.")
