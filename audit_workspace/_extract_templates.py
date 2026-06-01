"""Extract text from all docx + pptx files in AI Curriculum folder."""
from pathlib import Path
from docx import Document
from pptx import Presentation

CURR = Path("/Users/nalinaggarwal/Downloads/Supply-chain-main/AI Curriculum")

for path in sorted(CURR.iterdir()):
    if path.suffix == ".docx":
        print(f"\n\n{'='*80}\n=== {path.name}\n{'='*80}")
        doc = Document(path)
        for i, para in enumerate(doc.paragraphs):
            text = para.text
            if text.strip():
                style = para.style.name if para.style else ""
                print(f"[{style}] {text}")
        for ti, table in enumerate(doc.tables):
            print(f"\n--- TABLE {ti} ---")
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                print(" | ".join(cells))
    elif path.suffix == ".pptx":
        print(f"\n\n{'='*80}\n=== {path.name}\n{'='*80}")
        pres = Presentation(path)
        for si, slide in enumerate(pres.slides):
            print(f"\n--- SLIDE {si + 1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text.strip():
                            print(text)
                if shape.has_table:
                    print(" [TABLE]")
                    for row in shape.table.rows:
                        cells = [c.text_frame.text.strip() for c in row.cells]
                        print(" | ".join(cells))
